"""
Create PowerPoint Presentation from Dashboard Changes Summary
Requires: python-pptx library
Install: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_dashboard_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define IBM colors
    IBM_BLUE = RGBColor(0, 114, 206)
    IBM_DARK_BLUE = RGBColor(0, 67, 206)
    DARK_GRAY = RGBColor(50, 50, 50)
    LIGHT_GRAY = RGBColor(150, 150, 150)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Customer Ticketing Dashboard"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = IBM_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Comprehensive Improvements Summary"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "January - May 2026"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(20)
    date_para.font.color.rgb = LIGHT_GRAY
    date_para.alignment = PP_ALIGN.CENTER
    
    # Add presenter
    presenter_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "Prepared by: Jason Blakley"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(16)
    presenter_para.font.color.rgb = LIGHT_GRAY
    presenter_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "Transformed the dashboard into a scalable, enterprise-grade analytics platform"
    
    for bullet_text in [
        "15+ major improvements implemented",
        "70% performance improvement (10 min → 3 min load time)",
        "62% reduction in complexity (16 files → 6 files)",
        "Zero unplanned downtime maintained",
        "Comprehensive usage tracking at no additional cost",
        "28 months of historical data (Jan 2024 - Apr 2026)"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 3: The Challenge
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Challenge"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "Growing Pains"
    
    for bullet_text in [
        "Growing data volume (16+ monthly files)",
        "Deployment timeouts with large datasets",
        "No visibility into dashboard usage",
        "Complex and time-consuming monthly updates (60+ minutes)",
        "Performance degradation as data grew",
        "Risk of service disruption"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 4: Solution - Data Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Solution: Data Architecture Transformation"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "Consolidated Yearly File Structure"
    
    for bullet_text in [
        "Created Merged_data_2025.csv (684 MB)",
        "Reduced from 16 files to 6 files (62% reduction)",
        "70% faster startup time (10 min → 3 min)",
        "40% faster data loading",
        "Scalable for future growth",
        "Eliminated deployment timeout issues"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 5: Solution - Reliability & Monitoring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Solution: Reliability & Monitoring"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "Background Loading System"
    
    for bullet_text in [
        "Dual-port architecture (health check + main app)",
        "100% deployment success rate",
        "Passes health checks while loading 2+ GB of data",
        "",
        "Comprehensive Usage Tracking",
        "Track user logins and dashboard interactions",
        "Monitor internal vs external users",
        "Zero additional cost (uses existing logs)"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0 if bullet_text and not bullet_text.startswith("Track") and not bullet_text.startswith("Monitor") and not bullet_text.startswith("Zero") else 1
        p.font.size = Pt(18)
    
    # Slide 6: Performance Improvements
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Performance Improvements"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    # Add table
    rows, cols = 7, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(1.5)
    
    # Header row
    headers = ['Metric', 'Before', 'After', 'Improvement']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = IBM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    data = [
        ['File Count', '16 files', '6 files', '62% ↓'],
        ['Download Time', '3-5 min', '1-2 min', '60% ↓'],
        ['Load Time', '3-5 min', '2-3 min', '40% ↓'],
        ['Total Startup', '6-10 min', '3-5 min', '70% ↓'],
        ['Memory Usage', '16 DFs', '6 DFs', 'Major ↓'],
        ['Deployment', '80%', '100%', '20% ↑']
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Slide 7: Business Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Business Impact"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "For Sales Team"
    
    for bullet_text in [
        "70% faster dashboard access",
        "Current data (through April 2026)",
        "Improved accuracy with better date filtering",
        "Zero downtime during updates",
        "",
        "For External Clients",
        "Professional, fast experience",
        "Always up-to-date data",
        "Secure IBM App ID authentication",
        "",
        "For IT/Operations",
        "75% faster monthly updates (60 min → 15 min)",
        "Comprehensive usage monitoring",
        "Scalable for 3+ years of growth"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0 if bullet_text in ["For Sales Team", "For External Clients", "For IT/Operations"] else 1
        p.font.size = Pt(16)
        if bullet_text in ["For Sales Team", "For External Clients", "For IT/Operations"]:
            p.font.bold = True
    
    # Slide 8: Results & ROI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results & ROI"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "Key Metrics"
    
    for bullet_text in [
        "2.8+ million ticket records analyzed",
        "28 months of historical data coverage",
        "50+ major enterprise clients tracked",
        "2-3 active users daily",
        "",
        "Cost Efficiency",
        "Infrastructure: ~$55/month (unchanged)",
        "Development: 30 hours over 4 months",
        "Time saved: 45 min/month = 9 hours/year",
        "Avoided monitoring costs: $100+/month",
        "Zero downtime = maintained sales productivity"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0 if bullet_text in ["Key Metrics", "Cost Efficiency"] else 1
        p.font.size = Pt(16)
        if bullet_text in ["Key Metrics", "Cost Efficiency"]:
            p.font.bold = True
    
    # Slide 9: Future Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Roadmap"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "Short-term (Next Month)"
    
    for bullet_text in [
        "Complete April 2026 deployment",
        "Add May 2026 data (15 minutes)",
        "Implement automated data validation",
        "",
        "Medium-term (Next Quarter)",
        "Add 'Last Updated' indicator",
        "Implement quarterly consolidated files",
        "Add automated testing",
        "Create user feedback mechanism",
        "",
        "Long-term (Next 6 Months)",
        "Automated monthly data pipeline",
        "Predictive analytics features",
        "Mobile-responsive design",
        "Role-based access control"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0 if bullet_text in ["Short-term (Next Month)", "Medium-term (Next Quarter)", "Long-term (Next 6 Months)"] else 1
        p.font.size = Pt(16)
        if bullet_text in ["Short-term (Next Month)", "Medium-term (Next Quarter)", "Long-term (Next 6 Months)"]:
            p.font.bold = True
    
    # Slide 10: Summary & Questions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE
    
    content = slide.placeholders[1].text_frame
    content.text = "What We Achieved"
    
    for bullet_text in [
        "✓ Scalable, enterprise-grade analytics platform",
        "✓ 70% performance improvement",
        "✓ 62% reduction in complexity",
        "✓ Zero downtime maintained",
        "✓ Comprehensive usage tracking",
        "✓ 75% faster monthly updates",
        "",
        "Dashboard URL:",
        "https://python-appid-app.wt1yl0ero9k.us-south.codeengine.appdomain.cloud/dashboard/",
        "",
        "Questions?"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(18)
        if "Questions?" in bullet_text:
            p.font.bold = True
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = 'Dashboard_Changes_Presentation.pptx'
    prs.save(output_file)
    print(f"✓ Presentation created: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"\nOpen the file in PowerPoint to view and edit!")

if __name__ == "__main__":
    try:
        create_dashboard_presentation()
    except ImportError:
        print("Error: python-pptx library not found")
        print("\nPlease install it first:")
        print("  pip install python-pptx")
        print("\nThen run this script again:")
        print("  python create_presentation.py")

# Made with Bob
